#!/usr/bin/env python3
"""Extract text content from a .pptx and write a simple text-based PDF.

This does NOT try to visually render slides. It creates a readable PDF containing
slide-by-slide text content (titles/bullets/paragraphs).

Usage:
  python3 tools/pptx_content_to_pdf.py <input.pptx> [output.pdf]
"""

import sys
from pathlib import Path


def ensure_deps():
    try:
        import pptx  # noqa: F401
        from reportlab.pdfgen import canvas  # noqa: F401
        from reportlab.lib.pagesizes import letter  # noqa: F401
    except Exception as e:
        print("Missing dependency:", e)
        print("Install with:")
        print("  python3 -m pip install python-pptx reportlab")
        sys.exit(2)


def extract_text(pptx_path: Path):
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides_out = []

    for idx, slide in enumerate(prs.slides, start=1):
        items = []
        for shape in slide.shapes:
            if not hasattr(shape, "has_text_frame"):
                continue
            if not shape.has_text_frame:
                continue

            # Collect paragraphs; preserve basic bullet indentation via leading spaces.
            for p in shape.text_frame.paragraphs:
                text = "".join(run.text for run in p.runs).strip()
                if not text:
                    continue
                indent = "  " * int(getattr(p, "level", 0) or 0)
                items.append(f"{indent}{text}")

        slides_out.append({"slide": idx, "lines": items})

    return slides_out


def write_pdf(slides, out_path: Path, title: str):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    c = canvas.Canvas(str(out_path), pagesize=letter)
    width, height = letter

    margin_x = 54  # 0.75 inch
    margin_y = 54
    line_h = 14

    def draw_wrapped_lines(lines, x, y, max_width, font_name="Helvetica", font_size=11):
        from reportlab.pdfbase.pdfmetrics import stringWidth

        c.setFont(font_name, font_size)
        out = []
        for line in lines:
            words = line.split(" ")
            cur = ""
            for w in words:
                nxt = (cur + " " + w).strip() if cur else w
                if stringWidth(nxt, font_name, font_size) <= max_width:
                    cur = nxt
                else:
                    if cur:
                        out.append(cur)
                    cur = w
            if cur:
                out.append(cur)
        yy = y
        for l in out:
            c.drawString(x, yy, l)
            yy -= line_h
        return yy

    # Cover page
    c.setFont("Helvetica-Bold", 18)
    c.drawString(margin_x, height - margin_y, title)
    c.setFont("Helvetica", 11)
    c.drawString(margin_x, height - margin_y - 24, "Extracted slide text content")
    c.showPage()

    for s in slides:
        y = height - margin_y
        c.setFont("Helvetica-Bold", 14)
        c.drawString(margin_x, y, f"Slide {s['slide']}")
        y -= 22

        if not s["lines"]:
            c.setFont("Helvetica-Oblique", 11)
            c.drawString(margin_x, y, "(No text found on this slide)")
            c.showPage()
            continue

        # Draw lines with wrapping; paginate if needed.
        max_width = width - 2 * margin_x
        lines = s["lines"]
        cursor = y

        # chunk writing with page breaks
        buf = []
        for line in lines:
            buf.append(line)
            # conservative page break check
            if cursor - (len(buf) * line_h) < margin_y:
                cursor = draw_wrapped_lines(buf, margin_x, cursor, max_width)
                c.showPage()
                cursor = height - margin_y
                buf = []
        if buf:
            cursor = draw_wrapped_lines(buf, margin_x, cursor, max_width)
            c.showPage()

    c.save()


def main():
    ensure_deps()

    if len(sys.argv) < 2:
        print("Usage: python3 tools/pptx_content_to_pdf.py <input.pptx> [output.pdf]")
        sys.exit(1)

    in_path = Path(sys.argv[1]).expanduser().resolve()
    if not in_path.exists():
        print(f"Input not found: {in_path}")
        sys.exit(1)

    out_path = (
        Path(sys.argv[2]).expanduser().resolve()
        if len(sys.argv) >= 3
        else in_path.with_suffix(".content.pdf")
    )

    slides = extract_text(in_path)
    write_pdf(slides, out_path, title=in_path.stem)

    # Also print a short extraction summary to stdout
    total_lines = sum(len(s["lines"]) for s in slides)
    print(f"Created: {out_path}")
    print(f"Slides: {len(slides)} | Text lines: {total_lines}")


if __name__ == "__main__":
    main()
