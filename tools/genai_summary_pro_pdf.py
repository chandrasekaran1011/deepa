#!/usr/bin/env python3
"""Create a professional 1-page summary PDF from a PPTX.

- Extracts slide text (best effort)
- Extracts embedded images from the PPTX zip (ppt/media/*)
- Writes a styled PDF with background, typography, and simple infographic cards

Usage:
  python3 tools/genai_summary_pro_pdf.py <input.pptx> [output.pdf]

Output is intended as an executive summary, not a slide-perfect rendering.
"""

import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional


def ensure_deps():
    try:
        from pptx import Presentation  # noqa: F401
        from reportlab.lib.pagesizes import A4  # noqa: F401
        from reportlab.pdfgen import canvas  # noqa: F401
        from reportlab.lib import colors  # noqa: F401
        from reportlab.lib.utils import ImageReader  # noqa: F401
    except Exception as e:
        print("Missing dependency:", e)
        print("Install with:")
        print("  python3 -m pip install python-pptx reportlab")
        sys.exit(2)


@dataclass
class Extracted:
    title: str
    bullets: List[str]
    themes: List[str]
    use_cases: List[str]
    risks: List[str]
    metrics: List[str]
    images: List[Path]


def extract_slide_text(pptx_path: Path) -> List[List[str]]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides = []
    for slide in prs.slides:
        lines = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                text = "".join(run.text for run in p.runs).strip()
                if text:
                    level = int(getattr(p, "level", 0) or 0)
                    prefix = "- " if level == 0 else "  - "
                    lines.append(prefix + text)
        slides.append(lines)
    return slides


def extract_images(pptx_path: Path, out_dir: Path, max_images: int = 6) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    images = []
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        media_files = [n for n in z.namelist() if n.startswith("ppt/media/")]
        # Keep deterministic order
        for name in sorted(media_files):
            if len(images) >= max_images:
                break
            data = z.read(name)
            ext = Path(name).suffix.lower()
            if ext not in {".png", ".jpg", ".jpeg"}:
                continue
            out_path = out_dir / Path(name).name
            out_path.write_bytes(data)
            images.append(out_path)
    return images


def pick_summary(extracted_lines: List[List[str]]) -> Extracted:
    # Heuristic summary builder: pulls first non-empty line as title, then buckets by keywords.
    all_lines = [l.replace("- ", "").strip() for slide in extracted_lines for l in slide]
    all_lines = [l for l in all_lines if l]

    title = all_lines[0] if all_lines else "Innovation in GenAI — Executive Summary"

    def bucket(keywords):
        out = []
        for l in all_lines:
            ll = l.lower()
            if any(k in ll for k in keywords):
                out.append(l)
        # de-dupe preserving order
        seen = set()
        uniq = []
        for x in out:
            if x.lower() in seen:
                continue
            seen.add(x.lower())
            uniq.append(x)
        return uniq[:6]

    use_cases = bucket(["use case", "application", "customer", "support", "sales", "marketing", "agent", "assistant", "coding", "dev", "automation"])
    risks = bucket(["risk", "privacy", "security", "halluc", "compliance", "bias", "governance", "ip", "copyright"])
    metrics = bucket(["kpi", "metric", "measure", "latency", "cost", "quality", "accuracy", "roi", "adoption"])
    themes = bucket(["strategy", "roadmap", "architecture", "platform", "data", "model", "evaluation", "mvp", "scale"])

    # fallback: take some generic bullets if buckets are empty
    bullets = []
    for l in all_lines[1:]:
        if len(bullets) >= 8:
            break
        if l not in bullets:
            bullets.append(l)

    if not use_cases:
        use_cases = bullets[:4]
    if not themes:
        themes = bullets[4:8] if len(bullets) > 4 else bullets[:4]

    return Extracted(
        title=title,
        bullets=bullets[:10],
        themes=themes[:5],
        use_cases=use_cases[:5],
        risks=risks[:5],
        metrics=metrics[:5],
        images=[],
    )


def draw_wrapped(c, text: str, x: float, y: float, w: float, font: str, size: int, leading: float):
    from reportlab.pdfbase.pdfmetrics import stringWidth

    c.setFont(font, size)
    words = text.split()
    line = ""
    yy = y
    for word in words:
        test = (line + " " + word).strip()
        if stringWidth(test, font, size) <= w:
            line = test
        else:
            if line:
                c.drawString(x, yy, line)
                yy -= leading
            line = word
    if line:
        c.drawString(x, yy, line)
        yy -= leading
    return yy


def write_pro_pdf(ex: Extracted, out_path: Path, images: List[Path]):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    c = canvas.Canvas(str(out_path), pagesize=A4)
    W, H = A4

    # Palette
    bg = colors.HexColor("#0B1220")
    card = colors.HexColor("#111B2E")
    card2 = colors.HexColor("#0F1A33")
    accent = colors.HexColor("#6EE7FF")
    accent2 = colors.HexColor("#A78BFA")
    text = colors.HexColor("#E8EEF9")
    muted = colors.HexColor("#B7C3D8")

    # Background
    c.setFillColor(bg)
    c.rect(0, 0, W, H, stroke=0, fill=1)

    # Header band
    c.setFillColor(colors.HexColor("#0E1730"))
    c.rect(0, H - 120, W, 120, stroke=0, fill=1)

    # Accent line
    c.setFillColor(accent)
    c.rect(48, H - 86, 64, 4, stroke=0, fill=1)

    c.setFillColor(text)
    c.setFont("Helvetica-Bold", 22)
    c.drawString(48, H - 64, "Innovation in GenAI")
    c.setFont("Helvetica", 11)
    c.setFillColor(muted)
    c.drawString(48, H - 92, "Executive summary (content derived from the PPTX)")

    # Layout grid
    margin = 40
    gap = 14
    content_top = H - 140
    content_bottom = 42

    col_w = (W - 2 * margin - gap) / 2
    left_x = margin
    right_x = margin + col_w + gap

    def card_box(x, y, w, h, title, accent_color=accent):
        c.setFillColor(card)
        c.roundRect(x, y, w, h, 12, stroke=0, fill=1)
        c.setFillColor(accent_color)
        c.roundRect(x, y + h - 34, w, 34, 12, stroke=0, fill=1)
        c.setFillColor(colors.HexColor("#07101F"))
        c.setFont("Helvetica-Bold", 12)
        c.drawString(x + 14, y + h - 22, title)

    # Cards sizes
    left_h1 = 210
    left_h2 = 220
    right_h1 = 250
    right_h2 = 180

    y_left1 = content_top - left_h1
    y_left2 = y_left1 - gap - left_h2
    y_right1 = content_top - right_h1
    y_right2 = y_right1 - gap - right_h2

    card_box(left_x, y_left1, col_w, left_h1, "Key themes", accent)
    card_box(left_x, y_left2, col_w, left_h2, "High-impact use cases", accent2)
    card_box(right_x, y_right1, col_w, right_h1, "Risks & governance", colors.HexColor("#FCA5A5"))
    card_box(right_x, y_right2, col_w, right_h2, "Success metrics", colors.HexColor("#86EFAC"))

    # Bullet rendering
    def draw_bullets(x, y, w, h, bullets):
        c.setFillColor(text)
        yy = y + h - 50
        c.setFont("Helvetica", 10.5)
        for b in bullets[:7]:
            if yy < y + 18:
                break
            yy = draw_wrapped(c, f"• {b}", x + 14, yy, w - 28, "Helvetica", 10.5, 14)

    draw_bullets(left_x, y_left1, col_w, left_h1, ex.themes)
    draw_bullets(left_x, y_left2, col_w, left_h2, ex.use_cases)
    draw_bullets(right_x, y_right1, col_w, right_h1, ex.risks)
    draw_bullets(right_x, y_right2, col_w, right_h2, ex.metrics)

    # Simple infographic: maturity pipeline
    pipe_y = content_bottom + 10
    c.setFillColor(card2)
    c.roundRect(margin, pipe_y, W - 2 * margin, 78, 14, stroke=0, fill=1)
    c.setFillColor(text)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(margin + 16, pipe_y + 54, "Delivery approach")

    steps = ["Discover", "Prototype", "Pilot", "Scale"]
    step_w = (W - 2 * margin - 40) / 4
    sx = margin + 16
    sy = pipe_y + 20
    for i, s in enumerate(steps):
        c.setFillColor(accent if i < 2 else colors.HexColor("#334155"))
        c.roundRect(sx + i * step_w, sy, step_w - 10, 22, 10, stroke=0, fill=1)
        c.setFillColor(colors.HexColor("#07101F"))
        c.setFont("Helvetica-Bold", 10)
        c.drawCentredString(sx + i * step_w + (step_w - 10) / 2, sy + 7, s)

    # Image strip (optional): place up to 3 images as thumbnails (if present in ppt/media)
    if images:
        thumb_y = H - 120 - 12
        x0 = W - margin - 180
        c.setFillColor(colors.HexColor("#0E1730"))
        c.roundRect(x0 - 8, H - 112, 188, 60, 12, stroke=0, fill=1)
        max_thumbs = min(3, len(images))
        for i in range(max_thumbs):
            try:
                img = ImageReader(str(images[i]))
                c.drawImage(img, x0 + i * 58, H - 106, width=52, height=52, mask='auto', preserveAspectRatio=True, anchor='c')
            except Exception:
                pass

    # Footer
    c.setFillColor(muted)
    c.setFont("Helvetica", 8.5)
    c.drawString(margin, 20, "Generated from PPTX text content • Not a slide-accurate render")

    c.save()


def main():
    ensure_deps()

    if len(sys.argv) < 2:
        print("Usage: python3 tools/genai_summary_pro_pdf.py <input.pptx> [output.pdf]")
        sys.exit(1)

    in_path = Path(sys.argv[1]).expanduser().resolve()
    if not in_path.exists():
        print(f"Input not found: {in_path}")
        sys.exit(1)

    out_path = (
        Path(sys.argv[2]).expanduser().resolve()
        if len(sys.argv) >= 3
        else in_path.with_suffix(".summary.pdf")
    )

    extracted_lines = extract_slide_text(in_path)
    img_dir = in_path.parent / ".deepa_assets" / in_path.stem
    images = extract_images(in_path, img_dir)

    ex = pick_summary(extracted_lines)
    write_pro_pdf(ex, out_path, images)

    print(f"Created: {out_path}")
    print(f"Images embedded (thumbnails): {min(3, len(images))} (extracted: {len(images)})")


if __name__ == "__main__":
    main()
